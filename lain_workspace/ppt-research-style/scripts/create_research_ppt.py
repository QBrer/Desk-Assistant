#!/usr/bin/env python
"""Create Chinese research-style PowerPoint decks from a JSON spec."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:
    from add_speaker_notes import add_speaker_notes
except Exception:
    add_speaker_notes = None


SLIDE_W = 13.333
SLIDE_H = 7.5
ASSET_DIR = Path(__file__).resolve().parents[1] / "assets"
HEADER_STRIP = ASSET_DIR / "header-strip.png"
FOOTER_STRIP = ASSET_DIR / "footer-strip.png"
CN_FONT = "微软雅黑"
EN_FONT = "Times New Roman"
TEXT = RGBColor(15, 17, 21)
BODY = RGBColor(34, 34, 34)
BLUE = RGBColor(0, 97, 170)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(192, 0, 0)
LIGHT_LINE = RGBColor(217, 217, 217)


def _rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _set_east_asian_font(run, typeface: str) -> None:
    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag == qn("a:ea"):
            r_pr.remove(child)
    ea = OxmlElement("a:ea")
    ea.set("typeface", typeface)
    r_pr.append(ea)


def _format_run(run, size: float, bold: bool = False, color: RGBColor = BODY, font: str = CN_FONT) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _set_east_asian_font(run, font)


def _clear_text_frame(shape) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 16,
    bold: bool = False,
    color: RGBColor = BODY,
    font: str = CN_FONT,
    align: Optional[str] = None,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _clear_text_frame(shape)
    lines = str(text or "").splitlines() or [""]
    for index, line in enumerate(lines):
        p = shape.text_frame.paragraphs[0] if index == 0 else shape.text_frame.add_paragraph()
        p.space_after = Pt(3)
        p.line_spacing = 1.05
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        _format_run(run, size=size, bold=bold, color=color, font=font)
    return shape


def add_title(slide, title: str) -> None:
    add_text(slide, title, 0.38, 0.06, 12.0, 0.72, 32, True, WHITE, CN_FONT)


def add_section_bar(slide, text: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.96), Inches(SLIDE_W), Inches(0.48))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.color.rgb = WHITE
    _clear_text_frame(bar)
    p = bar.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    _format_run(run, 20, True, WHITE, CN_FONT)


def add_caption(slide, text: str, x: float, y: float, w: float, h: float = 0.3) -> None:
    if text:
        add_text(slide, text, x, y, w, h, 12, False, TEXT, CN_FONT, "center")


def add_citation(slide, text: str) -> None:
    if text:
        font = EN_FONT if all(ord(ch) < 128 for ch in text[:80]) else CN_FONT
        add_text(slide, text, 0.12, 7.02, 13.05, 0.30, 9.5, False, WHITE, font)


def add_placeholder(slide, label: str, x: float, y: float, w: float, h: float) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(248, 249, 250)
    rect.line.color.rgb = LIGHT_LINE
    add_text(slide, label or "Missing figure", x + 0.08, y + h / 2 - 0.18, max(w - 0.16, 0.5), 0.36, 12, False, BODY, CN_FONT, "center")


def fit_image(slide, image_path: str, x: float, y: float, w: float, h: float, label: str = "") -> None:
    path = Path(image_path) if image_path else Path()
    if not image_path or not path.exists() or not path.is_file():
        add_placeholder(slide, label or str(image_path), x, y, w, h)
        return

    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
    else:
        draw_h = h
        draw_w = h * img_ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), width=Inches(draw_w), height=Inches(draw_h))


def image_items(data: Dict[str, Any]) -> List[Dict[str, str]]:
    items = data.get("images") or []
    out = []
    for item in items:
        if isinstance(item, str):
            out.append({"path": item, "caption": ""})
        elif isinstance(item, dict):
            out.append({"path": item.get("path", ""), "caption": item.get("caption", "")})
    return out


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    if HEADER_STRIP.exists():
        slide.shapes.add_picture(str(HEADER_STRIP), Inches(0), Inches(0.01), width=Inches(SLIDE_W), height=Inches(0.96))
    if FOOTER_STRIP.exists():
        slide.shapes.add_picture(str(FOOTER_STRIP), Inches(0), Inches(6.86), width=Inches(SLIDE_W), height=Inches(0.64))
    return slide


def layout_split(slide, data: Dict[str, Any]) -> None:
    add_title(slide, data.get("title", ""))
    add_text(slide, data.get("body", ""), 0.0, 1.02, 6.65, 5.95, 16, False, BODY)
    boxes = [(6.82, 0.96, 6.15, 3.05, 4.62), (6.82, 4.10, 6.15, 2.55, 6.72)]
    for item, (x, y, w, h, cy) in zip(image_items(data), boxes):
        fit_image(slide, item["path"], x, y, w, h, item["caption"])
        add_caption(slide, item["caption"], x, cy, w)
    add_citation(slide, data.get("citation", ""))


def layout_evidence_pairs(slide, data: Dict[str, Any]) -> None:
    add_title(slide, data.get("title", ""))
    images = image_items(data)
    texts = data.get("blocks") or []
    figure_boxes = [(0.0, 1.02, 5.8, 2.73, 3.82), (0.0, 4.23, 5.88, 2.73, 6.93)]
    text_boxes = [(5.86, 1.14, 7.42, 2.29), (5.80, 4.43, 7.55, 2.29)]
    for item, (x, y, w, h, cy) in zip(images, figure_boxes):
        fit_image(slide, item["path"], x, y, w, h, item["caption"])
        add_caption(slide, item["caption"], x, cy, w)
    for text, (x, y, w, h) in zip(texts, text_boxes):
        add_text(slide, text, x, y, w, h, 16, False, BODY)
    if data.get("body") and not texts:
        add_text(slide, data["body"], 5.86, 1.14, 7.42, 5.5, 16, False, BODY)
    add_citation(slide, data.get("citation", ""))


def layout_case(slide, data: Dict[str, Any]) -> None:
    add_title(slide, data.get("title", ""))
    if data.get("section"):
        add_section_bar(slide, data["section"])
        body_y = 1.52
    else:
        body_y = 1.08
    add_text(slide, data.get("body", ""), 0.08, body_y, 13.05, 1.20, 16, False, BODY)
    images = image_items(data)
    if len(images) <= 1:
        boxes = [(0.8, 2.65, 11.7, 3.85, 6.55)]
    elif len(images) == 2:
        boxes = [(0.85, 2.55, 5.95, 3.95, 6.55), (7.25, 2.55, 5.45, 3.95, 6.55)]
    else:
        boxes = [(0.35, 2.35, 4.1, 3.95, 6.43), (4.62, 2.35, 4.1, 3.95, 6.43), (8.89, 2.35, 4.1, 3.95, 6.43)]
    for item, (x, y, w, h, cy) in zip(images, boxes):
        fit_image(slide, item["path"], x, y, w, h, item["caption"])
        add_caption(slide, item["caption"], x, cy, w)
    add_citation(slide, data.get("citation", ""))


def layout_grid(slide, data: Dict[str, Any]) -> None:
    add_title(slide, data.get("title", ""))
    add_text(slide, data.get("body", ""), 0.40, 1.10, 12.80, 0.78, 16, False, BODY)
    images = image_items(data)
    cols = int(data.get("columns", 3))
    cols = max(2, min(cols, 4))
    gap = 0.22
    start_x, start_y = 0.35, 2.05
    box_w = (12.65 - (cols - 1) * gap) / cols
    box_h = 1.75 if cols >= 3 else 2.25
    for idx, item in enumerate(images[: cols * 2]):
        row, col = divmod(idx, cols)
        x = start_x + col * (box_w + gap)
        y = start_y + row * (box_h + 0.58)
        fit_image(slide, item["path"], x, y, box_w, box_h, item["caption"])
        add_caption(slide, item["caption"], x, y + box_h + 0.06, box_w)
    add_citation(slide, data.get("citation", ""))


def layout_text(slide, data: Dict[str, Any]) -> None:
    add_title(slide, data.get("title", ""))
    if data.get("section"):
        add_section_bar(slide, data["section"])
        y = 1.52
    else:
        y = 1.08
    add_text(slide, data.get("body", ""), 0.10, y, 12.95, 5.55, 16, False, BODY)
    add_citation(slide, data.get("citation", ""))


LAYOUTS = {
    "split": layout_split,
    "evidence_pairs": layout_evidence_pairs,
    "case": layout_case,
    "grid": layout_grid,
    "text": layout_text,
}


def build_deck(spec: Dict[str, Any], output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    for slide_data in spec.get("slides", []):
        slide = new_slide(prs)
        layout = slide_data.get("layout", "case")
        LAYOUTS.get(layout, layout_case)(slide, slide_data)
    notes = [str(s.get("notes") or s.get("speaker_notes") or "") for s in spec.get("slides", [])]
    if any(note.strip() for note in notes) and add_speaker_notes is not None:
        base = output.with_name(output.stem + "_base.pptx")
        prs.save(base)
        add_speaker_notes(base, notes, output)
    else:
        prs.save(output)


def sample_spec() -> Dict[str, Any]:
    return {
        "deck_title": "自然资源大模型研究进展",
        "slides": [
            {
                "layout": "case",
                "title": "生物资源大模型-典型案例",
                "section": "AlphaFold 3大模型",
                "body": "打破前代模型仅限于单一蛋白质折叠的预测局限，具备对蛋白质、DNA、RNA、小分子配体及离子等核心分子的三维结构和复杂相互作用的联合预测能力。",
                "images": [
                    {"path": "", "caption": "AF3 推理架构"},
                    {"path": "", "caption": "使用 AF3 预测的复合物结构"},
                ],
                "citation": "Abramson, J. et al. Accurate structure prediction of biomolecular interactions with AlphaFold 3. Nature 630, 493-500 (2024).",
            },
            {
                "layout": "split",
                "title": "AlphaEarth Foundations (AEF)简介",
                "body": "AlphaEarth Foundations 是面向地理空间场景的基础模型。它将多源异构观测压缩为连续、可比较的特征表达，使稀疏观测能够支持跨区域、跨时间的自然资源分析。",
                "images": [
                    {"path": "", "caption": "AlphaEarth Foundations模型总架构"},
                    {"path": "", "caption": "包含全球主要陆地和微小岛屿完整嵌入场"},
                ],
            },
        ],
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("spec", nargs="?", help="Input JSON spec. Omit when using --sample.")
    parser.add_argument("output", nargs="?", help="Output .pptx path.")
    parser.add_argument("--sample", action="store_true", help="Generate a small sample deck.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if args.sample:
        output = Path(args.spec or args.output or "sample.pptx")
        build_deck(sample_spec(), output)
        print(f"Wrote {output}")
        return
    if not args.spec or not args.output:
        raise SystemExit("Usage: create_research_ppt.py input.json output.pptx")
    with open(args.spec, "r", encoding="utf-8-sig") as f:
        spec = json.load(f)
    build_deck(spec, Path(args.output))
    print(f"Wrote {args.output}")


if __name__ == "__main__":
    main()
